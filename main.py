import shutil
import os
import tempfile
import fitz  # PyMuPDF
from pptx import Presentation
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List

from app.services.ai_engine import get_engine

app = FastAPI(title="TemanStudi AI Service (PDF & PPTX)")

# CORS: Agar JS Backend bisa akses
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Output Schema (Sesuai ERD Project)
class FlashcardItem(BaseModel):
    pertanyaan: str
    jawaban: str

class AIResponse(BaseModel):
    status: str
    pesan: str
    data: List[FlashcardItem]

# --- Helper Ekstraksi ---

def extract_from_pdf(path: str, start: int, end: int) -> str:
    doc = fitz.open(path)
    text = ""
    total = len(doc)
    s_idx, e_idx = max(0, start - 1), min(total, end)
    for i in range(s_idx, e_idx):
        text += doc.load_page(i).get_text() + "\n"
    doc.close()
    return text

def extract_from_pptx(path: str, start: int, end: int) -> str:
    prs = Presentation(path)
    text = ""
    total = len(prs.slides)
    s_idx, e_idx = max(0, start - 1), min(total, end)
    for i in range(s_idx, e_idx):
        slide = prs.slides[i]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + ". "
        text += "\n"
    return text

# --- Endpoint Utama ---

@app.on_event("startup")
async def startup_event():
    """Panaskan mesin (Load Model ke GPU) saat server nyala"""
    get_engine()

@app.post("/generate", response_model=AIResponse)
async def generate_endpoint(
    file: UploadFile = File(...),
    start_page: int = Form(1),
    end_page: int = Form(10)
):
    # 1. Cek Ekstensi
    filename = file.filename.lower()
    if not (filename.endswith(".pdf") or filename.endswith(".pptx")):
        raise HTTPException(400, "Format file harus .pdf atau .pptx")

    # 2. Simpan Sementara
    ext = ".pdf" if filename.endswith(".pdf") else ".pptx"
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        shutil.copyfileobj(file.file, tmp)
        tmp_path = tmp.name

    try:
        # 3. Ekstrak Teks Berdasarkan Format
        full_text = ""
        if filename.endswith(".pdf"):
            full_text = extract_from_pdf(tmp_path, start_page, end_page)
        else:
            full_text = extract_from_pptx(tmp_path, start_page, end_page)

        # 4. Validasi Isi Teks (Menangani Scan Gambar/Slide Kosong)
        if len(full_text.strip()) < 50:
            return AIResponse(status="error", pesan="Teks kosong/terlalu pendek (Mungkin Gambar?)", data=[])

        # 5. Panggil AI Engine
        engine = get_engine()
        results = engine.generate_flashcards(full_text)

        # 6. Mapping ke Format ERD
        final_data = [
            FlashcardItem(pertanyaan=r.get("question",""), jawaban=r.get("answer","")) 
            for r in results
        ]

        return AIResponse(status="success", pesan="OK", data=final_data)

    except Exception as e:
        print(f"ERROR SYSTEM: {e}")
        raise HTTPException(500, str(e))
    
    finally:
        # Hapus file temp (Cleanup) agar storage tidak penuh
        if os.path.exists(tmp_path): os.remove(tmp_path)

@app.get("/")
def health():
    return {"status": "AI Ready", "model": "Qwen 3 1.7B"}
